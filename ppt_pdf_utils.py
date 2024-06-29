from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import comtypes.client
import os
from datetime import datetime
from openpyxl import load_workbook
import win32api

# Adobe Acrobat 的正确路径
acrobat_path = r"C:\Program Files\Adobe\Acrobat DC\Acrobat\Acrobat.exe"

def cm_to_inches(cm):
    return cm * 0.3937

def pptx_to_pdf(pptx_path, pdf_path):
    # 此處添加你將PPT轉為PDF的程式碼
    # 確保使用絕對路徑
    pptx_path = os.path.abspath(pptx_path)
    pdf_path = os.path.abspath(pdf_path)

    # 打印文件路徑以便於檢查
    print(f"Converting PPTX to PDF:\nPPTX Path: {pptx_path}\nPDF Path: {pdf_path}")
    
    # 確保 PPTX 文件存在
    if not os.path.exists(pptx_path):
        print(f"Error: PPTX file does not exist: {pptx_path}")
        return
    
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    
    # 嘗試打開 PPTX 文件
    try:
        ppt = powerpoint.Presentations.Open(pptx_path)
        ppt.SaveAs(pdf_path, FileFormat=32)  # 32 代表 PDF 格式
        ppt.Close()
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
    finally:
        powerpoint.Quit()

def process_excel_to_pdf(date):
    # 此處添加你處理Excel到PDF的程式碼
    template_path = r"C:\Users\lin\Documents\program\toss\Landseed\Landseed Model2.pptx"

    folder_path = f"motiphysio{date.replace('/', '')}"
    excel_folder_path = os.path.join(folder_path, "Excel")
    ppt_folder_path = os.path.join(folder_path, "PPT")
    final_pdf_folder = os.path.join(folder_path, "Final_PDF")

    os.makedirs(ppt_folder_path, exist_ok=True)
    os.makedirs(final_pdf_folder, exist_ok=True)

    # 確保圖片路徑和其他資源路徑存在並正確
    adductor_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Adductor1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Adductor2_corrected.jpg"]
    addStretch_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\AddStretch1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\AddStretch2_corrected.jpg"]
    bigBall_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Bigball1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Bigball2_corrected.jpg"]
    bridge_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Bridge1_corrected.jpg",r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Bridge2_corrected.jpg"]
    calf_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Calf1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Calf2_corrected.jpg"]
    chinIn_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\ChinIn1_corrected.jpg",r"C:\Users\lin\Documents\program\toss\Landseed\Picture\ChinIn2_corrected.jpg"]
    clam_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Clam1_corrected.jpg",r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Clam2_corrected.jpg"]
    deadbug_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Deadbug1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Deadbug2_corrected.jpg"]
    glut_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Glut1_corrected.jpg",r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Glut2_corrected.jpg"]
    hamstring_pictues = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Hamstring1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Hamstring2_corrected.jpg"]
    lowback_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\LBstretch1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\LBstretch2_corrected.jpg"]
    lowerTrap_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\LowerTrap1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\LowerTrap2_corrected.jpg"]
    oblique_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\OB1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\OB2_corrected.jpg"]
    pec_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Pec1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Pec2_corrected.jpg"]
    row_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Row1_corrected.jpg",r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Row2_corrected.jpg"]
    squat_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Squat1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Squat2_corrected.jpg"]
    upperTrap_pictures = [r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Trap1_corrected.jpg", r"C:\Users\lin\Documents\program\toss\Landseed\Picture\Trap2_corrected.jpg"]

    image_mapping = {
        "側弓箭步": adductor_pictures,
        "內收肌群伸展": addStretch_pictures,
        "抱大球": bigBall_pictures,
        "橋式": bridge_pictures,
        "小腿伸展": calf_pictures,
        "靠牆縮下巴": chinIn_pictures,
        "蚌殼式": clam_pictures,
        "死蟲式": deadbug_pictures,
        "臀肌伸展": glut_pictures,
        "腿後肌伸展": hamstring_pictues,
        "下背伸展": lowback_pictures,
        "上背伸展": lowerTrap_pictures,
        "腹斜肌伸展": oblique_pictures,
        "胸大肌伸展": pec_pictures,
        "彈力帶划船": row_pictures,
        "深蹲": squat_pictures,
        "斜方肌伸展": upperTrap_pictures
    }

    for excel_file in os.listdir(excel_folder_path):
        if excel_file.endswith(".xlsx"):
            excel_path = os.path.join(excel_folder_path, excel_file)
            wb = load_workbook(excel_path)
            ws = wb["對應試算表"]

            name = ws["A2"].value 
            rank = ws["B2"].value 
            date = ws["C2"].value
            f2 = ws["F2"].value 
            f3 = ws["F3"].value
            f4 = ws["F4"].value
            e2 = ws["E2"].value 
            e3 = ws["E3"].value
            e4 = ws["E4"].value
            g2 = ws["G2"].value 
            g3 = ws["G3"].value
            g4 = ws["G4"].value
            h2 = ws["H2"].value 
            h3 = ws["H3"].value 
            h4 = ws["H4"].value 
            i2 = ws["I2"].value 
            i3 = ws["I3"].value 
            i4 = ws["I4"].value 
            j2 = ws["J2"].value 
            j3 = ws["J3"].value 
            j4 = ws["J4"].value 
            k2 = ws["K2"].value 
            k3 = ws["K3"].value 
            k4 = ws["K4"].value

            prs = Presentation(template_path)
            slide1 = prs.slides[0]
            slide2 = prs.slides[1]

            page_mapping = {
                "name": 0, "date": 0, "rank": 0, 
                "f2": 0, "e2": 0, "g2": 0, "h2": 0, "i2": 0, "j2": 0, "k2": 0, 
                "p1.1": 0, "p1.2": 0, "p1.3": 0, "p1.4": 0, 
                "f3": 1, "e3": 1, "g3": 1, "h3": 1, "i3": 1, "j3": 1, "k3": 1, 
                "p2.1": 1, "p2.2": 1, "p2.3": 1, "p2.4": 1,
                "f4": 1, "e4": 1, "g4": 1, "h4": 1, "i4": 1, "j4": 1, "k4": 1, 
                "p3.1": 1, "p3.2": 1, "p3.3": 1, "p3.4": 1
            }

            target_positions = {
                "name": (3.4, 4), 
                "date": (12.8, 4), 
                "rank": (1.6, 7.3),
                "f2": (4.2, 12.6), 
                "f3": (4.2, 1.5), 
                "f4": (4.2, 15), 
                "e2": (4.2, 14.3), 
                "e3": (4.2, 3.3), 
                "e4": (4.2, 16.8), 
                "g2": (4.2, 13.5), 
                "g3": (4.2, 2.4), 
                "g4": (4.2, 15.9), 
                "h2": (4.2, 15.3), 
                "h3": (4.2, 4.2), 
                "h4": (4.2, 17.6), 
                "i2": (8, 16.2), 
                "i3": (8, 5.1), 
                "i4": (8, 18.5), 
                "j2": (4.2, 19.5), 
                "j3": (4.2, 8.6), 
                "j4": (4.2, 22), 
                "k2": (8, 20.3), 
                "k3": (8, 9.6), 
                "k4": (8, 22.9), 
                "p1.1": (2.2, 16.2), 
                "p1.2": (5, 16.2),
                "p1.3": (2.2, 20.3), 
                "p1.4": (5, 20.3),
                "p2.1": (2.2, 5.1), 
                "p2.2": (5, 5.1), 
                "p2.3": (2.2, 9.6), 
                "p2.4": (5, 9.6),
                "p3.1": (2.2, 18.5),
                "p3.2": (5, 18.5),
                "p3.3": (2.2, 22.9),
                "p3.4": (5, 22.9),
            }

            tolerance = Inches(0.05)

            for slide_index, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for key, (x_cm, y_cm) in target_positions.items():
                            if page_mapping[key] == slide_index:
                                target_x = Inches(cm_to_inches(x_cm))
                                target_y = Inches(cm_to_inches(y_cm))

                                if abs(shape.left - target_x) <= tolerance and abs(shape.top - target_y) <= tolerance:
                                    text_frame = shape.text_frame
                                    paragraph = text_frame.paragraphs[0]
                                    value = locals().get(key)
                                    if value:
                                        paragraph.text = str(value)

                                    if not paragraph.runs:
                                        paragraph.add_run()

                                    font = paragraph.runs[0].font

                                    if key in ["name", "date"]:
                                        font.name = "Microsoft JhengHei"
                                        font.size = Pt(18)
                                        font.bold = True
                                        font.color.rgb = RGBColor(0, 112, 192)

                                    elif key == "rank":
                                        font.size = Pt(66)
                                        font.bold = True
                                        font.color.rgb = RGBColor(51, 153, 255)

                                    elif key in ["f2", "f3", "f4"]:
                                        font.name = "Microsoft JhengHei"
                                        font.size = Pt(14)
                                        font.bold = True
                                        font.color.rgb = RGBColor(0, 112, 192)

                                    elif key in ["e2", "e3", "e4", "g2", "g3", "g4"]:
                                        font.name = "Microsoft JhengHei"
                                        font.size = Pt(12)
                                        font.bold = False
                                        font.color.rgb = RGBColor(0, 0, 0)

                                    elif key in ["h2", "h3", "h4", "j2", "j3", "j4"]:
                                        font.name = "Microsoft JhengHei"
                                        font.size = Pt(12)
                                        font.bold = True
                                        font.color.rgb = RGBColor(0, 112, 192)

                                    elif key in ["i2", "i3", "i4", "k2", "k3", "k4"]:
                                        font.name = "Microsoft JhengHei"
                                        font.size = Pt(10)
                                        font.bold = False
                                        font.color.rgb = RGBColor(0, 0, 0)

            if h2 in image_mapping:
                images = image_mapping[h2]
                pos1 = target_positions["p1.1"]
                pos2 = target_positions["p1.2"]
                slide1.shapes.add_picture(images[0], Inches(cm_to_inches(pos1[0])), Inches(cm_to_inches(pos1[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))
                slide1.shapes.add_picture(images[1], Inches(cm_to_inches(pos2[0])), Inches(cm_to_inches(pos2[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))

            if j2 in image_mapping:
                images = image_mapping[j2]
                pos1 = target_positions["p1.3"]
                pos2 = target_positions["p1.4"]
                slide1.shapes.add_picture(images[0], Inches(cm_to_inches(pos1[0])), Inches(cm_to_inches(pos1[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))
                slide1.shapes.add_picture(images[1], Inches(cm_to_inches(pos2[0])), Inches(cm_to_inches(pos2[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))

            if h3 in image_mapping:
                images = image_mapping[h3]
                pos1 = target_positions["p2.1"]
                pos2 = target_positions["p2.2"]
                slide2.shapes.add_picture(images[0], Inches(cm_to_inches(pos1[0])), Inches(cm_to_inches(pos1[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))
                slide2.shapes.add_picture(images[1], Inches(cm_to_inches(pos2[0])), Inches(cm_to_inches(pos2[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))

            if j3 in image_mapping:
                images = image_mapping[j3]
                pos1 = target_positions["p2.3"]
                pos2 = target_positions["p2.4"]
                slide2.shapes.add_picture(images[0], Inches(cm_to_inches(pos1[0])), Inches(cm_to_inches(pos1[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))
                slide2.shapes.add_picture(images[1], Inches(cm_to_inches(pos2[0])), Inches(cm_to_inches(pos2[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))

            if h4 in image_mapping:
                images = image_mapping[h4]
                pos1 = target_positions["p3.1"]
                pos2 = target_positions["p3.2"]
                slide2.shapes.add_picture(images[0], Inches(cm_to_inches(pos1[0])), Inches(cm_to_inches(pos1[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))
                slide2.shapes.add_picture(images[1], Inches(cm_to_inches(pos2[0])), Inches(cm_to_inches(pos2[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))

            if j4 in image_mapping:
                images = image_mapping[j4]
                pos1 = target_positions["p3.3"]
                pos2 = target_positions["p3.4"]
                slide2.shapes.add_picture(images[0], Inches(cm_to_inches(pos1[0])), Inches(cm_to_inches(pos1[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))
                slide2.shapes.add_picture(images[1], Inches(cm_to_inches(pos2[0])), Inches(cm_to_inches(pos2[1])), width=Inches(cm_to_inches(2.43)), height=Inches(cm_to_inches(3.25)))

            base_name = os.path.splitext(excel_file)[0].split("_")[0] + "_moti中文報告"
            pptx_file = f"{base_name}.pptx"
            pptx_path = os.path.join(ppt_folder_path, pptx_file)
            prs.save(pptx_path)

            pdf_file = f"{base_name}.pdf"
            pdf_path = os.path.join(final_pdf_folder, pdf_file)
            pptx_to_pdf(pptx_path, pdf_path)

def print_pdf(file_to_print):
    # 此處添加你列印PDF的程式碼
    if os.path.exists(file_to_print):
        if os.path.exists(acrobat_path):
            try:
                win32api.ShellExecute(
                    0,
                    "open",
                    acrobat_path,
                    f'/p /h "{file_to_print}"',
                    ".",
                    0
                )
                print("文件已发送至默认打印机")
            except Exception as e:
                print(f"列印失败: {e}")
        else:
            print("找不到 Adobe Acrobat 的执行档，请确认安装路径")
    else:
        print("文件不存在，请检查文件路径")
